import os
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("A biblioteca 'python-pptx' não está instalada.")
    print("Por favor, execute: pip install python-pptx")
    exit(1)

def criar_apresentacao():
    prs = Presentation()
    
    # Slide 1: Capa
    slide_layout = prs.slide_layouts[0] # Layout de Título
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "O caos na gestão acaba aqui."
    subtitle.text = "Projeto Athena - Plataforma SaaS de Previsibilidade e Controle"
    
    # Slide 2: O Problema
    slide_layout = prs.slide_layouts[1] # Título e Conteúdo
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "O Desafio"
    tf = content.text_frame
    tf.text = "Projetos atrasam porque os riscos não são vistos a tempo."
    
    p = tf.add_paragraph()
    p.text = "Diretores e gestores não precisam de mais um 'to-do list'."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Eles precisam de uma visão executiva clara de onde os gargalos estão ocorrendo."
    p2.level = 1
    
    # Slide 3: A Solução
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "A Solução - Visão 360º"
    tf = content.text_frame
    tf.text = "Dashboard Executivo"
    
    p = tf.add_paragraph()
    p.text = "Performance, etapas e orçamentos centralizados."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Decisões baseadas em dados reais de consumo e avanço."
    p2.level = 1
    
    # Slide 4: Inovação
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Inovação - Risk Map"
    tf = content.text_frame
    tf.text = "Evite o burnout e atrase de entregas"
    
    p = tf.add_paragraph()
    p.text = "O Mapa de Risco identifica automaticamente membros da equipe sobrecarregados."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Permite o balanceamento de carga de trabalho antes do projeto sofrer impactos."
    p2.level = 1
    
    # Slide 5: Conclusão
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    
    title.text = "Construído com excelência"
    tf = content.text_frame
    tf.text = "O que você achou da interface?"
    
    p = tf.add_paragraph()
    p.text = "O Athena foi desenvolvido buscando a melhor experiência visual e técnica (React/TypeScript)."
    p.level = 1
    
    p2 = tf.add_paragraph()
    p2.text = "Deixe seu feedback nos comentários!"
    p2.level = 1

    # Salva o arquivo
    arquivo_saida = "Apresentacao_Athena.pptx"
    prs.save(arquivo_saida)
    print(f"Sucesso! Arquivo '{arquivo_saida}' gerado na pasta atual.")

if __name__ == "__main__":
    criar_apresentacao()
